"""Generate MemeSonic presentation slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

BASE = os.path.dirname(os.path.abspath(__file__))
RESULTS = os.path.join(BASE, "results")
OUT = os.path.join(BASE, "MemeSonic_slides.pptx")

# ── colour palette ────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x0D, 0x1A)   # near-black
ACCENT   = RGBColor(0x7C, 0x3A, 0xED)   # purple
ACCENT2  = RGBColor(0x06, 0xB6, 0xD4)   # cyan
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xCC, 0xCC, 0xCC)
DGRAY    = RGBColor(0x33, 0x33, 0x55)
GREEN    = RGBColor(0x10, 0xB9, 0x81)
YELLOW   = RGBColor(0xF5, 0x9E, 0x0B)
RED      = RGBColor(0xEF, 0x44, 0x44)

SW, SH = Inches(13.33), Inches(7.5)   # 16:9 widescreen

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

blank = prs.slide_layouts[6]   # completely blank layout

# ── helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    s = prs.slides.add_slide(blank)
    bg = s.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return s

def tb(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def img(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)

def hline(slide, y, color=DGRAY, w=None):
    line_w = w or SW
    shp = slide.shapes.add_shape(1, Inches(0), y, line_w, Pt(1))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()

def chip(slide, text, l, t, w=Inches(1.8), h=Inches(0.38), color=ACCENT):
    r = rect(slide, l, t, w, h, color)
    tb(slide, text, l, t, w, h, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def section_tag(slide, text):
    rect(slide, Inches(0), Inches(0), Inches(2.6), Inches(0.42), ACCENT)
    tb(slide, text, Inches(0.1), Inches(0.04), Inches(2.5), Inches(0.38),
       size=13, bold=True, color=WHITE)

def title_slide_heading(slide, title, sub=None):
    tb(slide, title, Inches(1.2), Inches(2.4), Inches(10.9), Inches(1.4),
       size=44, bold=True, color=WHITE)
    if sub:
        tb(slide, sub, Inches(1.2), Inches(3.9), Inches(10.9), Inches(0.6),
           size=20, color=LGRAY)

def slide_title(slide, text):
    tb(slide, text, Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.7),
       size=30, bold=True, color=WHITE)
    hline(slide, Inches(1.35), color=ACCENT, w=SW)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, Inches(0), Inches(0), Inches(0.25), SH, ACCENT)          # left stripe
rect(s, Inches(0), SH - Inches(0.25), SW, Inches(0.25), ACCENT2) # bottom stripe

tb(s, "MemeSonic", Inches(1.0), Inches(1.6), Inches(11), Inches(1.2),
   size=52, bold=True, color=WHITE)
tb(s, "Multimodal Meme Understanding for Emotion-Aware Audio Generation",
   Inches(1.0), Inches(2.9), Inches(11), Inches(0.8), size=22, color=ACCENT2)
tb(s, "Image-Text Fusion · Incongruity Modeling · FusionMoE",
   Inches(1.0), Inches(3.8), Inches(11), Inches(0.5), size=16, color=LGRAY)
tb(s, "MAS.60 Final Project  ·  Spring 2026",
   Inches(1.0), Inches(6.5), Inches(6), Inches(0.5), size=14, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Project Overview & Pipeline
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Project Overview")

tb(s, "Goal: Detect meme sentiment & intention → generate emotion-matched audio",
   Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.5), size=16, color=LGRAY)

# pipeline boxes
boxes = [
    ("Meme Input\n(image + text)", DGRAY),
    ("Multimodal\nFusion", ACCENT),
    ("Sentiment &\nIntention", ACCENT2),
    ("Audio\nGeneration", GREEN),
]
bw, bh = Inches(2.3), Inches(1.35)
by = Inches(2.3)
for i, (label, col) in enumerate(boxes):
    bx = Inches(0.6) + i * Inches(3.1)
    rect(s, bx, by, bw, bh, col)
    tb(s, label, bx, by, bw, bh, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(boxes) - 1:
        tb(s, "→", bx + bw, by + Inches(0.45), Inches(0.8), Inches(0.45),
           size=22, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# key tasks
tasks = [
    ("Sentiment", "7-class emotion\n(happiness, sorrow, …)"),
    ("Intention", "4-class purpose\n(expressive, entertaining, …)"),
    ("Sarcasm", "Binary + 3-class\nacross datasets"),
    ("Humor", "3-class\n(not funny → hilarious)"),
]
ty = Inches(4.0)
for i, (name, desc) in enumerate(tasks):
    tx = Inches(0.5) + i * Inches(3.2)
    chip(s, name, tx, ty, Inches(1.5), Inches(0.35), ACCENT)
    tb(s, desc, tx, ty + Inches(0.42), Inches(2.8), Inches(0.7),
       size=13, color=LGRAY)

tb(s, "3 Datasets:  MET-Meme (3,389) · Memotion-7k (6,992) · MMSD 2.0",
   Inches(0.5), Inches(5.5), Inches(12), Inches(0.4), size=14, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Dataset EDA
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "01  Datasets")
slide_title(s, "Meme Datasets: What We're Working With")

col_stats = [
    ("MET-Meme", ["3,389 samples", "7-class sentiment", "5-class intention", "Hate · Metaphor tasks"]),
    ("Memotion-7k", ["6,992 samples", "Humor · Sarcasm", "Offensive · Motivational", "Strong class imbalance"]),
    ("MMSD 2.0", ["2,409 test samples", "Binary sarcasm detection", "Cleaner labels", "Benchmark dataset"]),
]
for i, (title, bullets) in enumerate(col_stats):
    cx = Inches(0.5) + i * Inches(4.3)
    rect(s, cx, Inches(1.6), Inches(4.0), Inches(3.8), DGRAY)
    tb(s, title, cx + Inches(0.15), Inches(1.7), Inches(3.7), Inches(0.5),
       size=18, bold=True, color=ACCENT2)
    for j, b in enumerate(bullets):
        tb(s, f"• {b}", cx + Inches(0.15), Inches(2.3) + j * Inches(0.52),
           Inches(3.7), Inches(0.45), size=14, color=WHITE)

# imbalance callout
rect(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.55), RGBColor(0x1E, 0x1E, 0x3A))
tb(s, "⚠  Class imbalance: positive sentiment 59.5% · funny 67% · motivational 65% not-motivational  →  use Macro-F1 as primary metric",
   Inches(0.65), Inches(5.75), Inches(12), Inches(0.45), size=13, color=YELLOW)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — LLM Baseline
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "02  LLM Baseline")
slide_title(s, "LLM Baseline: GPT-4o & Gemini-2.5-Flash")

# result table
headers = ["Task", "GPT-4o", "Gemini-2.5"]
rows = [
    ("MET-Meme Intention",  "52%",  "48%",  "medium"),
    ("MET-Meme Sentiment",  "28%",  "34%",  "bad"),
    ("MET-Meme Hate",       "76%",  "74%",  "good"),
    ("MET-Meme Metaphor",   "58%",  "72%",  "medium"),
    ("Memotion Humor",      "68%",  "68%",  "medium"),
    ("Memotion Sarcasm",    "52%",  "50%",  "medium"),
    ("Memotion Sentiment",  "30%",  "26%",  "bad"),
    ("MMSD Sarcasm",        "80%",  "86%",  "good"),
]
col_colors = {"good": GREEN, "medium": YELLOW, "bad": RED}
col_w = [Inches(3.4), Inches(1.5), Inches(1.5)]
col_x = [Inches(0.5), Inches(4.0), Inches(5.6)]
row_h = Inches(0.45)
hy = Inches(1.55)

# header row
for ci, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    rect(s, cx, hy, cw, row_h, ACCENT)
    tb(s, hdr, cx, hy, cw, row_h, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, (task, g4, gem, rating) in enumerate(rows):
    ry = hy + row_h + ri * row_h
    bg_col = RGBColor(0x1A, 0x1A, 0x2E) if ri % 2 == 0 else DGRAY
    for ci, (cx, cw) in enumerate(zip(col_x, col_w)):
        rect(s, cx, ry, cw, row_h, bg_col)
    val_color = col_colors[rating]
    tb(s, task,  col_x[0], ry, col_w[0], row_h, size=13, color=WHITE, align=PP_ALIGN.LEFT)
    tb(s, g4,    col_x[1], ry, col_w[1], row_h, size=13, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    tb(s, gem,   col_x[2], ry, col_w[2], row_h, size=13, bold=True, color=val_color, align=PP_ALIGN.CENTER)

# takeaways
takeaways = [
    ("Hate / binary tasks", "relatively easy (74–86%)", GREEN),
    ("Sentiment & Intention", "very hard (26–52%)", RED),
    ("Interactive intention", "0% recall across ALL models", RED),
    ("Label subjectivity", "limits ceiling performance", YELLOW),
]
tx = Inches(7.4)
tb(s, "Key Takeaways", tx, Inches(1.55), Inches(5.4), Inches(0.45),
   size=16, bold=True, color=ACCENT2)
for i, (label, val, col) in enumerate(takeaways):
    ty = Inches(2.1) + i * Inches(0.82)
    rect(s, tx, ty, Inches(5.4), Inches(0.72), DGRAY)
    tb(s, label, tx + Inches(0.1), ty + Inches(0.04), Inches(5.1), Inches(0.3),
       size=13, bold=True, color=col)
    tb(s, val,   tx + Inches(0.1), ty + Inches(0.34), Inches(5.1), Inches(0.3),
       size=12, color=LGRAY)

tb(s, "→  LLMs alone are insufficient for fine-grained meme understanding",
   Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4), size=14,
   bold=True, color=ACCENT2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — What Is Incongruity
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "03  Incongruity Fusion")
slide_title(s, "The Key Insight: Meme Humor Lives in Incongruity")

tb(s, "A meme's meaning often emerges from the gap between image and text — not from either alone.",
   Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.5), size=16, color=LGRAY)

# three types
types = [
    ("Redundancy",  "Image & text say\nthe same thing",   DGRAY,   "→ simple concat works"),
    ("Incongruity", "Image & text\ncontradict each other", ACCENT,  "→ capture the tension"),
    ("Synergy",     "Meaning only emerges\nfrom both together", ACCENT2, "→ joint reasoning"),
]
bw, bh = Inches(3.6), Inches(2.2)
by = Inches(2.2)
for i, (name, desc, col, note) in enumerate(types):
    bx = Inches(0.5) + i * Inches(4.2)
    rect(s, bx, by, bw, bh, col)
    tb(s, name, bx + Inches(0.15), by + Inches(0.15), bw - Inches(0.3), Inches(0.5),
       size=18, bold=True, color=WHITE)
    tb(s, desc, bx + Inches(0.15), by + Inches(0.75), bw - Inches(0.3), Inches(0.8),
       size=14, color=WHITE)
    tb(s, note, bx + Inches(0.15), by + Inches(1.65), bw - Inches(0.3), Inches(0.4),
       size=12, color=LGRAY)

# architecture description
rect(s, Inches(0.5), Inches(4.65), Inches(12.3), Inches(1.55), DGRAY)
tb(s, "Incongruity-Aware Fusion Architecture",
   Inches(0.65), Inches(4.72), Inches(12), Inches(0.4), size=15, bold=True, color=ACCENT2)
arch_steps = [
    "① Image → ResNet-50 → projection",
    "② Text → BERT [CLS] → projection",
    "③ Both → sentiment distribution space",
    "④ Cross-modal disagreement = incongruity score",
    "⑤ Fused representation for classification",
]
for i, step in enumerate(arch_steps):
    col_ = i // 3
    row_ = i % 3
    tb(s, step,
       Inches(0.65) + col_ * Inches(6.0),
       Inches(5.15) + row_ * Inches(0.32),
       Inches(5.8), Inches(0.3), size=13, color=WHITE)

tb(s, "Trainable params: only 200,336  (frozen BERT + ResNet backbones)",
   Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.35), size=13, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Incongruity Fusion Results
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "03  Incongruity Fusion")
slide_title(s, "Incongruity Fusion Results: Wins & Limits")

img_path = os.path.join(RESULTS, "metaphor_model_comparison.png")
img(s, img_path, Inches(0.4), Inches(1.5), Inches(6.2), Inches(4.2))

# results summary on right
results = [
    ("Metaphor",   "81% → 87.7%", "+6.5 pp", True),
    ("MMSD Sarcasm","82.7% → 84.1%", "+1.4 pp", True),
    ("Hate",       "83.5% → 87.1%", "+3.6 pp", True),
    ("Sentiment",  "43.8% → 47.9%", "+4.1 pp", True),
    ("Intention",  "42.4% → 43.5%", "+1.1 pp", False),
]
tb(s, "Concat baseline  →  Incongruity Fusion",
   Inches(6.9), Inches(1.55), Inches(6.0), Inches(0.4), size=14, bold=True, color=LGRAY)

for i, (task, change, delta, good) in enumerate(results):
    ry = Inches(2.05) + i * Inches(0.82)
    col_ = GREEN if good else YELLOW
    rect(s, Inches(6.9), ry, Inches(5.9), Inches(0.72), DGRAY)
    tb(s, task,   Inches(7.05), ry + Inches(0.04), Inches(2.5), Inches(0.3), size=14, bold=True, color=WHITE)
    tb(s, change, Inches(7.05), ry + Inches(0.36), Inches(3.5), Inches(0.28), size=12, color=LGRAY)
    tb(s, delta,  Inches(11.2), ry + Inches(0.18), Inches(1.4), Inches(0.35),
       size=14, bold=True, color=col_, align=PP_ALIGN.RIGHT)

rect(s, Inches(6.9), Inches(6.2), Inches(5.9), Inches(0.6), RGBColor(0x1E, 0x1E, 0x3A))
tb(s, "✓ Strong on incongruity-heavy tasks (metaphor, sarcasm)\n✗ Limited gains on fine-grained sentiment & intention",
   Inches(7.05), Inches(6.22), Inches(5.6), Inches(0.55), size=12, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Why We Need MoE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "04  FusionMoE")
slide_title(s, "The Problem: One Fusion Strategy Doesn't Fit All Tasks")

tb(s, "Different tasks require different fusion strategies — a single architecture can't optimise for all of them.",
   Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.45), size=15, color=LGRAY)

evidence = [
    ("Metaphor", "High incongruity\nbetween image & text", "Incongruity fusion\nworks best", ACCENT),
    ("Sentiment", "Both modalities\ncarry unique signals", "Late fusion\nworks better", ACCENT2),
    ("Sarcasm\n(MMSD)", "Balanced contribution\nfrom both modalities", "No clear\nwinner fusion", GREEN),
    ("Intention", "Complex & imbalanced\nclass distribution", "All fusion\nstrategies struggle", RED),
]
bw, bh = Inches(2.9), Inches(2.6)
by = Inches(2.15)
for i, (task, prob, sol, col) in enumerate(evidence):
    bx = Inches(0.4) + i * Inches(3.2)
    rect(s, bx, by, bw, bh, DGRAY)
    rect(s, bx, by, bw, Inches(0.08), col)   # color top bar
    tb(s, task, bx + Inches(0.12), by + Inches(0.15), bw - Inches(0.2), Inches(0.55),
       size=15, bold=True, color=col)
    tb(s, prob, bx + Inches(0.12), by + Inches(0.75), bw - Inches(0.2), Inches(0.75),
       size=13, color=LGRAY)
    hline(s, by + Inches(1.6), color=RGBColor(0x44, 0x44, 0x66), w=bw)
    tb(s, sol,  bx + Inches(0.12), by + Inches(1.72), bw - Inches(0.2), Inches(0.75),
       size=13, bold=True, color=WHITE)

rect(s, Inches(0.4), Inches(5.05), Inches(12.5), Inches(0.65), ACCENT)
tb(s, "Solution: Let the model learn WHICH fusion strategy to apply, conditioned on the incongruity signal  →  FusionMoE",
   Inches(0.55), Inches(5.1), Inches(12.2), Inches(0.55), size=15, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FusionMoE Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "04  FusionMoE")
slide_title(s, "FusionMoE Architecture")

tb(s, "Three experts, each capturing a different fusion mode. An incongruity-aware router assigns weights.",
   Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.45), size=15, color=LGRAY)

experts = [
    ("Expert R\n(Redundancy)", "Concat Fusion\nShared info from both", ACCENT),
    ("Expert U\n(Uniqueness)",  "Late Fusion\nEach modality's\nunique contribution", ACCENT2),
    ("Expert S\n(Synergy)",     "Incongruity Fusion\nCaptures cross-modal\ntension", GREEN),
]
ew, eh = Inches(3.0), Inches(2.4)
ey = Inches(2.2)
for i, (name, desc, col) in enumerate(experts):
    ex = Inches(0.5) + i * Inches(3.5)
    rect(s, ex, ey, ew, eh, DGRAY)
    rect(s, ex, ey, ew, Inches(0.08), col)
    tb(s, name, ex + Inches(0.12), ey + Inches(0.12), ew - Inches(0.2), Inches(0.65),
       size=15, bold=True, color=col)
    tb(s, desc, ex + Inches(0.12), ey + Inches(0.85), ew - Inches(0.2), Inches(1.0),
       size=13, color=LGRAY)

# router box
rect(s, Inches(10.7), Inches(2.2), Inches(2.2), Inches(2.4), RGBColor(0x1E, 0x1E, 0x3A))
tb(s, "Router", Inches(10.82), Inches(2.32), Inches(2.0), Inches(0.4),
   size=14, bold=True, color=ACCENT2)
tb(s, "Incongruity\nscore → softmax\ngate weights\n(learned temp τ)",
   Inches(10.82), Inches(2.75), Inches(2.0), Inches(1.5), size=12, color=LGRAY)

# arrow down to output
tb(s, "↓  weighted sum  →  classifier",
   Inches(0.5), Inches(4.75), Inches(12.3), Inches(0.4),
   size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# backbone note
rect(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.55), DGRAY)
tb(s, "Backbone: Pre-cached CLIP embeddings (512-dim, L2-normalized)  ·  Trained end-to-end on each task",
   Inches(0.65), Inches(5.25), Inches(12.0), Inches(0.45), size=13, color=LGRAY)

# gate weight preview
rect(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.2), RGBColor(0x12, 0x12, 0x28))
tb(s, "Gate weight examples (mean across test set):",
   Inches(0.65), Inches(5.98), Inches(12.0), Inches(0.35), size=13, bold=True, color=ACCENT2)
gate_examples = [
    ("Metaphor",  "Concat 0.87 · LateFusion 0.13"),
    ("Sentiment", "Concat 0.44 · LateFusion 0.56"),
    ("Sarcasm",   "Concat 0.50 · LateFusion 0.50"),
    ("Hate",      "Concat 0.49 · LateFusion 0.51"),
]
for i, (task, weights) in enumerate(gate_examples):
    tx = Inches(0.65) + i * Inches(3.2)
    tb(s, task,    tx, Inches(6.38), Inches(3.0), Inches(0.3), size=13, bold=True, color=WHITE)
    tb(s, weights, tx, Inches(6.70), Inches(3.0), Inches(0.3), size=11, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FusionMoE Results (with gate weight heatmap)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "04  FusionMoE")
slide_title(s, "FusionMoE Results")

img(s, os.path.join(RESULTS, "moe_gate_weights.png"),
    Inches(0.4), Inches(1.5), Inches(6.0), Inches(4.0))

# result table right side
tb(s, "Key results (Accuracy · Macro-F1)",
   Inches(6.7), Inches(1.55), Inches(6.2), Inches(0.4), size=15, bold=True, color=LGRAY)

result_rows = [
    ("MET-Meme Metaphor",  "87.65%", "87.65%", True),
    ("MET-Meme Hate",      "87.06%", "38.26%", True),
    ("MET-Meme Sentiment", "48.24%", "26.81%", True),
    ("MMSD Sarcasm",       "84.10%", "83.97%", True),
    ("Memotion Humor",     "65.86%", "28.28%", True),
    ("MET-Meme Intention", "40.00%", "27.31%", False),
]
hdr_cols = ["Task", "Acc", "F1"]
hdr_xs   = [Inches(6.7), Inches(10.6), Inches(11.8)]
hdr_ws   = [Inches(3.8), Inches(1.1),  Inches(1.1)]
rh = Inches(0.42)
hy2 = Inches(2.0)
for ci, (h, hx, hw) in enumerate(zip(hdr_cols, hdr_xs, hdr_ws)):
    rect(s, hx, hy2, hw, rh, ACCENT)
    tb(s, h, hx, hy2, hw, rh, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for ri, (task, acc, f1, good) in enumerate(result_rows):
    ry = hy2 + rh + ri * rh
    bg = RGBColor(0x1A, 0x1A, 0x2E) if ri % 2 == 0 else DGRAY
    col_ = GREEN if good else YELLOW
    for hx, hw in zip(hdr_xs, hdr_ws):
        rect(s, hx, ry, hw, rh, bg)
    tb(s, task, hdr_xs[0], ry, hdr_ws[0], rh, size=12, color=WHITE)
    tb(s, acc,  hdr_xs[1], ry, hdr_ws[1], rh, size=12, bold=True,
       color=col_, align=PP_ALIGN.CENTER)
    tb(s, f1,   hdr_xs[2], ry, hdr_ws[2], rh, size=12,
       color=LGRAY, align=PP_ALIGN.CENTER)

tb(s, "← Gate weights show the router learns task-specific fusion strategies",
   Inches(0.4), Inches(5.6), Inches(6.2), Inches(0.4), size=12, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Full Heatmap
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "04  FusionMoE")
slide_title(s, "Model Comparison: All Tasks × All Methods")

img(s, os.path.join(RESULTS, "moe_full_heatmap.png"),
    Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))

tb(s, "Darker = higher accuracy.  FusionMoE / MMOE consistently competitive; incongruity fusion strong on metaphor & sarcasm.",
   Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35), size=12, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Error Analysis & Limitations
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "05  Analysis")
slide_title(s, "Error Analysis & Limitations")

items = [
    (RED,    "Interactive intention: 0% recall across ALL models",
             "Rarest class, semantically ambiguous — models default to majority classes"),
    (YELLOW, "Macro-F1 << Accuracy on imbalanced tasks",
             "Sentiment positive=60%, funny=67% → high accuracy from majority class voting"),
    (YELLOW, "Incongruity fusion hurts intention",
             "42.4% (concat) → 36.5% (incongruity) — cross-modal disagreement is noisy for this task"),
    (GREEN,  "FusionMoE gate learns task signal",
             "Metaphor →  concat-heavy  |  Sentiment → late-fusion-heavy  (without any supervision)"),
    (ACCENT2,"LLaVA probing (supplementary)",
             "Cross-modal attention differs between sarcastic / non-sarcastic at specific layers & heads"),
]
for i, (col, title, body) in enumerate(items):
    ry = Inches(1.6) + i * Inches(0.98)
    rect(s, Inches(0.4), ry, Inches(12.5), Inches(0.88), DGRAY)
    rect(s, Inches(0.4), ry, Inches(0.12), Inches(0.88), col)
    tb(s, title, Inches(0.65), ry + Inches(0.05), Inches(12.1), Inches(0.35),
       size=14, bold=True, color=col)
    tb(s, body,  Inches(0.65), ry + Inches(0.42), Inches(12.1), Inches(0.35),
       size=13, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Summary & Next Steps
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Summary & Next Steps")

rect(s, Inches(0), Inches(0), Inches(0.25), SH, ACCENT)

contrib = [
    "LLM baseline → fine-grained meme tasks are hard (26–52% on sentiment/intention)",
    "Incongruity-Aware Fusion → captures cross-modal tension, strong on metaphor (+6.5 pp) & sarcasm",
    "FusionMoE → task-adaptive routing; router learns fusion preference without supervision",
    "Systematic analysis across 10 tasks, 3 datasets, 8+ model architectures",
]
tb(s, "Contributions", Inches(0.5), Inches(1.55), Inches(7.5), Inches(0.4),
   size=18, bold=True, color=ACCENT2)
for i, c in enumerate(contrib):
    tb(s, f"✓  {c}", Inches(0.5), Inches(2.05) + i * Inches(0.58),
       Inches(7.5), Inches(0.5), size=14, color=WHITE)

hline(s, Inches(4.5), color=DGRAY, w=Inches(7.8))

next_steps = [
    "Connect FusionMoE output to audio generation (tone / mood conditioning)",
    "Address Interactive class: data augmentation or few-shot prompting",
    "Scale to larger CLIP / LLM backbone for richer embeddings",
]
tb(s, "Next Steps", Inches(0.5), Inches(4.6), Inches(7.5), Inches(0.4),
   size=18, bold=True, color=ACCENT2)
for i, n in enumerate(next_steps):
    tb(s, f"→  {n}", Inches(0.5), Inches(5.1) + i * Inches(0.52),
       Inches(7.5), Inches(0.45), size=14, color=LGRAY)

# right panel — metric highlight
rect(s, Inches(8.5), Inches(1.5), Inches(4.4), Inches(5.6), DGRAY)
tb(s, "Best Results", Inches(8.65), Inches(1.6), Inches(4.1), Inches(0.4),
   size=16, bold=True, color=ACCENT2)
highlights = [
    ("Metaphor",      "87.65%", "Acc"),
    ("MMSD Sarcasm",  "84.10%", "Acc"),
    ("Hate",          "87.06%", "Acc"),
    ("Humor",         "65.86%", "Acc"),
    ("Motivational",  "52.9%",  "F1 (LLM-MoE)"),
]
for i, (task, val, metric) in enumerate(highlights):
    ty = Inches(2.1) + i * Inches(0.85)
    tb(s, task,   Inches(8.65), ty,               Inches(2.5), Inches(0.35), size=13, color=LGRAY)
    tb(s, val,    Inches(8.65), ty + Inches(0.35), Inches(2.5), Inches(0.38),
       size=22, bold=True, color=GREEN)
    tb(s, metric, Inches(11.1), ty + Inches(0.42), Inches(1.6), Inches(0.3), size=11, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Saved → {OUT}")
